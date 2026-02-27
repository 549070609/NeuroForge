#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成简洁橘黄色科技风PPT
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

def set_slide_background(slide, r, g, b):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)

def add_text_box(slide, left, top, width, height, text, font_size=24, bold=False, color=(255, 255, 255), align=PP_ALIGN.LEFT):
    """添加文本框"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(color[0], color[1], color[2])
    p.alignment = align
    return textbox

def add_bullet_text(slide, left, top, width, height, items, font_size=20, color=(255, 255, 255)):
    """添加项目符号文本"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(color[0], color[1], color[2])
        p.space_before = Pt(10)
        p.space_after = Pt(10)
    return textbox

def create_ppt():
    """创建PPT"""
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # 定义配色方案 - 橘黄色科技风
    DARK_BG = (20, 20, 30)  # 深色背景
    ORANGE_PRIMARY = (255, 140, 0)  # 橘黄色主色
    ORANGE_LIGHT = (255, 180, 50)  # 浅橘黄色
    TEXT_WHITE = (255, 255, 255)  # 白色文本
    TEXT_GRAY = (200, 200, 200)  # 灰色文本

    # ========== 1. 封面 ==========
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, *DARK_BG)

    # 主标题
    add_text_box(slide, Inches(1), Inches(3), Inches(14), Inches(1.5),
                 "Vibe Coding 实践分享", font_size=60, bold=True,
                 color=ORANGE_PRIMARY, align=PP_ALIGN.CENTER)

    # 副标题
    add_text_box(slide, Inches(1), Inches(4.5), Inches(14), Inches(1),
                 "AI 时代的编程新范式", font_size=32,
                 color=ORANGE_LIGHT, align=PP_ALIGN.CENTER)

    # ========== 2. 总览介绍 ==========
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, *DARK_BG)

    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(15), Inches(1),
                 "总览介绍", font_size=48, bold=True,
                 color=ORANGE_PRIMARY, align=PP_ALIGN.LEFT)

    items = [
        "dev Spark - 开发工具平台",
        "Neuro Forge \"神经熔炉\" - AI Agent框架",
        "BMAD - 开发方法论"
    ]
    add_bullet_text(slide, Inches(1), Inches(2), Inches(14), Inches(6),
                    items, font_size=28, color=TEXT_WHITE)

    # ========== 3. Vibe Coding 过程介绍 ==========
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, *DARK_BG)

    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(15), Inches(1),
                 "Vibe Coding 过程介绍", font_size=48, bold=True,
                 color=ORANGE_PRIMARY, align=PP_ALIGN.LEFT)

    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(15), Inches(0.8),
                 "忘掉代码的存在", font_size=28, bold=True,
                 color=ORANGE_LIGHT, align=PP_ALIGN.LEFT)

    items = [
        "一个灵感",
        "灵感扩展（头脑风暴）",
        "构建文档",
        "寻找可借鉴开源项目（Learn code）",
        "总结借鉴项目",
        "构建项目文档",
        "实现核心版本",
        "快速迭代扩展"
    ]
    add_bullet_text(slide, Inches(1), Inches(2.5), Inches(14), Inches(6),
                    items, font_size=24, color=TEXT_WHITE)

    # ========== 4. dev Spark ==========
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, *DARK_BG)

    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(15), Inches(1),
                 "dev Spark", font_size=48, bold=True,
                 color=ORANGE_PRIMARY, align=PP_ALIGN.LEFT)

    add_text_box(slide, Inches(1), Inches(2), Inches(14), Inches(1),
                 "在汽车（vibe coding）代替马车（古法编程）的时代，做自动驾驶平台",
                 font_size=28, color=ORANGE_LIGHT, align=PP_ALIGN.LEFT)

    add_text_box(slide, Inches(1), Inches(3.5), Inches(14), Inches(1),
                 "核心功能：AI 驱动的开发工具平台",
                 font_size=24, color=TEXT_WHITE, align=PP_ALIGN.LEFT)

    # ========== 5. Neuro Forge ==========
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, *DARK_BG)

    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(15), Inches(1),
                 "Neuro Forge \"神经熔炉\"", font_size=48, bold=True,
                 color=ORANGE_PRIMARY, align=PP_ALIGN.LEFT)

    add_text_box(slide, Inches(1), Inches(2), Inches(14), Inches(0.8),
                 "核心架构", font_size=32, bold=True,
                 color=ORANGE_LIGHT, align=PP_ALIGN.LEFT)

    items = [
        "感知器系统",
        "双日志 + 本体模式",
        "模块化架构设计",
        "Agent 协作框架"
    ]
    add_bullet_text(slide, Inches(1), Inches(3), Inches(14), Inches(5),
                    items, font_size=24, color=TEXT_WHITE)

    # ========== 6. BMAD ==========
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, *DARK_BG)

    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(15), Inches(1),
                 "BMAD - 开发与产品互补短板", font_size=48, bold=True,
                 color=ORANGE_PRIMARY, align=PP_ALIGN.LEFT)

    items = [
        "产品需求文档 (PRD) 生成",
        "用户体验设计",
        "架构设计",
        "敏捷开发流程",
        "AI Agent 协作开发"
    ]
    add_bullet_text(slide, Inches(1), Inches(2), Inches(14), Inches(6),
                    items, font_size=24, color=TEXT_WHITE)

    # ========== 7. 反思 ==========
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, *DARK_BG)

    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(15), Inches(1),
                 "反思", font_size=48, bold=True,
                 color=ORANGE_PRIMARY, align=PP_ALIGN.LEFT)

    items = [
        "认知决定论：从A到B最难的不是如何走，而是要先知道有B可以走",
        "生态位重新定义：程序员 + 产品经理 + 业务专家 三位一体",
        "开发模式重新定义：介于敏捷与不敏捷之间",
        "不要站在马夫的视角看汽车：关注发展，关注收益"
    ]
    add_bullet_text(slide, Inches(1), Inches(2), Inches(14), Inches(6),
                    items, font_size=22, color=TEXT_WHITE)

    # ========== 8. 问题 ==========
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, *DARK_BG)

    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(15), Inches(1),
                 "问题", font_size=48, bold=True,
                 color=ORANGE_PRIMARY, align=PP_ALIGN.LEFT)

    items = [
        "信息爆炸：AI反馈海量信息，信息超负荷",
        "认知盲区：存在认知盲区，无法驱动工具",
        "质量控制：AI输出无法一次完美，需要质量控制",
        "人成瓶颈：人的精力、设计、需求明确度跟不上AI"
    ]
    add_bullet_text(slide, Inches(1), Inches(2), Inches(14), Inches(6),
                    items, font_size=24, color=TEXT_WHITE)

    # ========== 9. 结束页 ==========
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, *DARK_BG)

    add_text_box(slide, Inches(1), Inches(3.5), Inches(14), Inches(2),
                 "谢谢观看", font_size=60, bold=True,
                 color=ORANGE_PRIMARY, align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(5.5), Inches(14), Inches(1),
                 "Questions & Discussion", font_size=32,
                 color=ORANGE_LIGHT, align=PP_ALIGN.CENTER)

    # 保存PPT
    output_path = "vibe_coding_presentation.pptx"
    prs.save(output_path)
    print(f"PPT已生成: {os.path.abspath(output_path)}")

if __name__ == "__main__":
    try:
        create_ppt()
    except Exception as e:
        print(f"生成PPT时出错: {e}")
        import traceback
        traceback.print_exc()
